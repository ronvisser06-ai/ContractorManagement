"""Feature 2, Step 3 — proves the deterministic extractor against two cases
the brief calls out explicitly: (1) slides/text/tables/speaker-notes fidelity,
and (2) a linked-but-missing video surfacing as a warning, not a crash.

No real customer deck was available for case (1) when this was first
written, so a synthetic fixture covers it deterministically; a real 63-slide
deck (SampleOrientation/) is exercised separately as a smoke test.
"""

import io
import os

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches
from pypdf import PdfWriter

import extract_pdf
import extract_pptx

extract_pptx.upload_object = lambda *a, **k: None  # no Storage in unit tests
extract_pdf.upload_object = lambda *a, **k: None


def _tiny_png() -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (1, 1), color="white").save(buf, format="PNG")
    return buf.getvalue()


def _build_fixture_pptx() -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # "Title and Content"

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Confined Space Entry"
    body = slide.placeholders[1].text_frame
    body.text = "Never enter without a valid permit."
    p2 = body.add_paragraph()
    p2.text = "Atmosphere must be tested before entry."

    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(3), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "Hazard"
    table_shape.table.cell(0, 1).text = "Control"
    table_shape.table.cell(1, 0).text = "Oxygen deficiency"
    table_shape.table.cell(1, 1).text = "Continuous gas monitoring"

    slide.notes_slide.notes_text_frame.text = "Emphasize that permits expire after each shift."

    video_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    pic = video_slide.shapes.add_picture(io.BytesIO(_tiny_png()), Inches(1), Inches(1), Inches(2), Inches(2))
    rid = video_slide.part.relate_to(
        "file:///C:/decks/confined_space.mp4", RT.VIDEO, is_external=True
    )
    nv_pr = pic._element.find(f"{qn('p:nvPicPr')}/{qn('p:nvPr')}")
    video_file_el = etree.SubElement(nv_pr, qn("a:videoFile"))
    video_file_el.set(qn("r:link"), rid)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_fixture_slides_tables_notes_and_linked_video():
    deck = extract_pptx.extract_pptx(_build_fixture_pptx(), job_id="job_test", site_id="site_test")

    assert deck["source"]["type"] == "pptx"
    assert deck["source"]["slide_count"] == 2

    slide0 = deck["slides"][0]
    assert slide0["title"] == "Confined Space Entry"
    texts = [tr["text"] for tr in slide0["text_runs"]]
    assert "Never enter without a valid permit." in texts
    assert "Atmosphere must be tested before entry." in texts
    assert slide0["speaker_notes"] == "Emphasize that permits expire after each shift."

    assert len(slide0["tables"]) == 1
    table = slide0["tables"][0]
    assert table["headers"] == ["Hazard", "Control"]
    assert table["rows"] == [["Oxygen deficiency", "Continuous gas monitoring"]]

    slide1 = deck["slides"][1]
    assert len(slide1["media_asset_ids"]) == 1
    video_asset_id = slide1["media_asset_ids"][0]
    video_asset = next(a for a in deck["assets"] if a["id"] == video_asset_id)
    assert video_asset["embed_state"] == "linked_missing"
    assert video_asset["source_uri"] == "file:///C:/decks/confined_space.mp4"

    assert deck["warnings"] == [
        {"code": "MEDIA_LINK_UNRESOLVED", "slide_index": 1, "asset_id": video_asset_id}
    ]


_SAMPLE_DECK_PATH = os.path.join(
    os.path.dirname(__file__), "..", "SampleOrientation", "2025 Proton Safety Orientation_V2.0_Draft (1).pptx"
)


def test_extract_real_deck_smoke():
    if not os.path.exists(_SAMPLE_DECK_PATH):
        return  # sample deck isn't committed to the repo; skip outside this machine
    with open(_SAMPLE_DECK_PATH, "rb") as f:
        data = f.read()

    deck = extract_pptx.extract_pptx(data, job_id="job_smoke", site_id="site_smoke")

    assert deck["source"]["slide_count"] == 63
    assert len(deck["slides"]) == 63
    # Every slide carries at least text or an image — nothing was silently dropped.
    assert all(s["text_runs"] or s["image_asset_ids"] for s in deck["slides"])
    assert deck["warnings"] == []


def test_extract_pdf_does_not_crash():
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)

    deck = extract_pdf.extract_pdf(buf.getvalue(), job_id="job_test", site_id="site_test")

    assert deck["source"]["type"] == "pdf"
    assert deck["source"]["slide_count"] == 1
    assert deck["slides"][0]["id"] == "slide_0"
    assert deck["warnings"] == []
