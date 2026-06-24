"""Deterministic PPTX -> ExtractedDeck parse (contracts §4.1). No LLM calls.

Asset handling: embedded images/video are uploaded to Storage here and
referenced by storage_key + sha256 (never inlined). A video or picture whose
relationship target_mode is External (i.e. it points at a file outside the
package, the .pptx "linked media" case) is recorded with
embed_state="linked_missing" and a MEDIA_LINK_UNRESOLVED warning instead of
being treated as a fatal error — the deck still extracts.
"""

import hashlib
import io
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from storage import upload_object

_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _iter_shapes(shapes):
    """Flatten group shapes so nested shapes are visited too."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _find_media_rid(shape) -> str | None:
    """An <a:videoFile>/<a:audioFile> element anywhere under this shape's XML
    carries the media relationship id on r:link (linked) or r:embed
    (embedded), regardless of whether python-pptx classifies the shape as
    PICTURE (poster image + linked video) or MEDIA."""
    for el in shape._element.iter():
        if etree.QName(el.tag).localname in ("videoFile", "audioFile", "wavAudioFile"):
            return el.get(f"{_REL_NS}link") or el.get(f"{_REL_NS}embed")
    return None


def _find_blip_rid(shape) -> str | None:
    """python-pptx's shape.image only reads the main <a:blip r:embed|r:link>
    attribute. Modern Office vector icons store the reference instead on an
    <a16:svgBlip>/<asvg:svgBlip> extension inside the same blip's <a:extLst>
    — still a perfectly normal embedded package part, just not one
    shape.image knows to look at. Fall back to any r:embed/r:link found
    anywhere under the shape so those aren't misreported as linked-missing."""
    for el in shape._element.iter():
        rid = el.get(f"{_REL_NS}embed") or el.get(f"{_REL_NS}link")
        if rid:
            return rid
    return None


def _paragraph_bold(paragraph) -> bool:
    return any(run.font.bold for run in paragraph.runs if run.font.bold is not None)


def _slide_text_runs(slide) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    text_runs: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []

    for shape_index, shape in enumerate(_iter_shapes(slide.shapes)):
        if shape.has_table:
            table = shape.table
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            headers, *body_rows = rows if rows else ([], [])
            tables.append({"shape_index": shape_index, "headers": headers, "rows": body_rows})
            continue

        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            text_runs.append(
                {
                    "shape_index": shape_index,
                    "level": paragraph.level or 0,
                    "bold": _paragraph_bold(paragraph),
                    "text": text,
                }
            )

    return text_runs, tables


def _asset_id(slide_index: int, counter: int) -> str:
    return f"ast_{slide_index}_{counter}"


def _slide_media_assets(
    slide, slide_index: int, site_id: str, job_id: str, counter: list[int]
) -> tuple[list[dict[str, Any]], list[str], list[str], list[dict[str, Any]]]:
    """Returns (assets, image_asset_ids, media_asset_ids, warnings) for one slide."""
    assets: list[dict[str, Any]] = []
    image_ids: list[str] = []
    media_ids: list[str] = []
    warnings: list[dict[str, Any]] = []

    for shape in _iter_shapes(slide.shapes):
        rid = _find_media_rid(shape)
        if rid is not None:
            rel = shape.part.rels[rid]
            asset_id = _asset_id(slide_index, counter[0])
            counter[0] += 1
            if rel.is_external:
                assets.append(
                    {
                        "id": asset_id,
                        "kind": "video",
                        "storage_key": None,
                        "mime": None,
                        "embed_state": "linked_missing",
                        "source_uri": rel.target_ref,
                    }
                )
                warnings.append(
                    {"code": "MEDIA_LINK_UNRESOLVED", "slide_index": slide_index, "asset_id": asset_id}
                )
            else:
                part = rel.target_part
                blob = part.blob
                storage_key = f"sites/{site_id}/jobs/{job_id}/assets/{asset_id}"
                upload_object(storage_key, blob, part.content_type)
                assets.append(
                    {
                        "id": asset_id,
                        "kind": "video",
                        "storage_key": storage_key,
                        "mime": part.content_type,
                        "sha256": hashlib.sha256(blob).hexdigest(),
                        "embed_state": "embedded",
                    }
                )
            media_ids.append(asset_id)
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            asset_id = _asset_id(slide_index, counter[0])
            counter[0] += 1
            try:
                image = shape.image
                blob, content_type, ext, width, height = (
                    image.blob,
                    image.content_type,
                    image.ext,
                    *image.size,
                )
            except Exception:
                # shape.image only follows the main <a:blip> embed/link — fall back to
                # any blip extension (e.g. a modern SVG icon) before assuming the
                # picture is genuinely linked-and-missing.
                rid = _find_blip_rid(shape)
                rel = shape.part.rels[rid] if rid else None
                if rel is not None and not rel.is_external:
                    part = rel.target_part
                    blob, content_type, ext, width, height = (
                        part.blob,
                        part.content_type,
                        part.partname.ext.lstrip("."),
                        None,
                        None,
                    )
                else:
                    assets.append(
                        {
                            "id": asset_id,
                            "kind": "image",
                            "storage_key": None,
                            "mime": None,
                            "embed_state": "linked_missing",
                            "source_uri": rel.target_ref if rel is not None else None,
                        }
                    )
                    warnings.append(
                        {"code": "MEDIA_LINK_UNRESOLVED", "slide_index": slide_index, "asset_id": asset_id}
                    )
                    image_ids.append(asset_id)
                    continue

            storage_key = f"sites/{site_id}/jobs/{job_id}/assets/{asset_id}.{ext}"
            upload_object(storage_key, blob, content_type)
            assets.append(
                {
                    "id": asset_id,
                    "kind": "image",
                    "storage_key": storage_key,
                    "mime": content_type,
                    "sha256": hashlib.sha256(blob).hexdigest(),
                    "width": width,
                    "height": height,
                    "embed_state": "embedded",
                    "alt": shape.name,
                }
            )
            image_ids.append(asset_id)

    return assets, image_ids, media_ids, warnings


def _branding(prs: Presentation) -> dict[str, Any]:
    """Best-effort theme read. A miss here doesn't fail the job — branding is
    cosmetic and the contract field is left null when not confidently found."""
    branding: dict[str, Any] = {
        "colors": {"primary": None, "secondary": None, "accent": None},
        "fonts": {"heading": None, "body": None},
        "logo_asset_id": None,
    }
    try:
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        theme_part = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        theme_el = etree.fromstring(theme_part.blob)
        color_scheme = theme_el.find(".//a:clrScheme", ns)
        if color_scheme is not None:
            srgb = color_scheme.findall(".//a:srgbClr", ns)
            slots = ["primary", "secondary", "accent"]
            for slot, node in zip(slots, srgb[:3]):
                branding["colors"][slot] = f"#{node.get('val')}"
        font_scheme = theme_el.find(".//a:fontScheme", ns)
        if font_scheme is not None:
            major = font_scheme.find(".//a:majorFont/a:latin", ns)
            minor = font_scheme.find(".//a:minorFont/a:latin", ns)
            if major is not None:
                branding["fonts"]["heading"] = major.get("typeface")
            if minor is not None:
                branding["fonts"]["body"] = minor.get("typeface")
    except Exception:
        pass
    return branding


def extract_pptx(data: bytes, *, job_id: str, site_id: str) -> dict[str, Any]:
    prs = Presentation(io.BytesIO(data))

    all_assets: list[dict[str, Any]] = []
    all_warnings: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []
    counter = [1]

    for slide_index, slide in enumerate(prs.slides):
        text_runs, tables = _slide_text_runs(slide)
        assets, image_ids, media_ids, warnings = _slide_media_assets(
            slide, slide_index, site_id, job_id, counter
        )
        all_assets.extend(assets)
        all_warnings.extend(warnings)

        title = None
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip() or None

        speaker_notes = ""
        if slide.has_notes_slide:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(
            {
                "index": slide_index,
                "id": f"slide_{slide_index}",
                "title": title,
                "text_runs": text_runs,
                "tables": tables,
                "image_asset_ids": image_ids,
                "media_asset_ids": media_ids,
                "speaker_notes": speaker_notes,
            }
        )

    return {
        "source": {
            "type": "pptx",
            "slide_count": len(prs.slides),
            "sha256": hashlib.sha256(data).hexdigest(),
        },
        "branding": _branding(prs),
        "assets": all_assets,
        "slides": slides,
        "warnings": all_warnings,
    }
